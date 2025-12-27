#!/usr/bin/env python3
"""
Create a single-slide executive presentation for MedFlow
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_executive_slide():
    """Create a single impactful slide for executive management"""

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add blank slide
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Define colors
    dark_blue = RGBColor(0, 32, 96)
    medium_blue = RGBColor(0, 102, 204)
    green = RGBColor(0, 128, 0)
    white = RGBColor(255, 255, 255)
    light_gray = RGBColor(245, 245, 245)

    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = white

    # ===== HEADER =====
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "MedFlow: All-in-One Healthcare Platform"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = dark_blue

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Modernize Healthcare Delivery | Optimize Revenue | Enhance Patient Experience"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = medium_blue

    # ===== ROI HIGHLIGHT BOX =====
    roi_left = Inches(0.5)
    roi_top = Inches(1.8)
    roi_width = Inches(9)
    roi_height = Inches(1.2)

    # ROI background box
    roi_shape = slide.shapes.add_shape(
        1,  # Rectangle
        roi_left, roi_top, roi_width, roi_height
    )
    roi_shape.fill.solid()
    roi_shape.fill.fore_color.rgb = RGBColor(0, 102, 204)
    roi_shape.line.color.rgb = RGBColor(0, 102, 204)

    # ROI Title
    roi_title_box = slide.shapes.add_textbox(roi_left + Inches(0.2), roi_top + Inches(0.15), roi_width - Inches(0.4), Inches(0.3))
    roi_title_frame = roi_title_box.text_frame
    roi_title_frame.text = "Financial Impact - Year 1"
    roi_title_para = roi_title_frame.paragraphs[0]
    roi_title_para.alignment = PP_ALIGN.CENTER
    roi_title_para.font.size = Pt(20)
    roi_title_para.font.bold = True
    roi_title_para.font.color.rgb = white

    # ROI Metrics (3 columns)
    metric_top = roi_top + Inches(0.5)
    col_width = Inches(2.8)

    # Metric 1: ROI
    m1_box = slide.shapes.add_textbox(roi_left + Inches(0.3), metric_top, col_width, Inches(0.6))
    m1_frame = m1_box.text_frame
    m1_para = m1_frame.paragraphs[0]
    m1_para.alignment = PP_ALIGN.CENTER

    run1 = m1_para.add_run()
    run1.text = "400-856%"
    run1.font.size = Pt(32)
    run1.font.bold = True
    run1.font.color.rgb = white

    run2 = m1_para.add_run()
    run2.text = "\nROI"
    run2.font.size = Pt(14)
    run2.font.color.rgb = white

    # Metric 2: Payback
    m2_box = slide.shapes.add_textbox(roi_left + Inches(3.3), metric_top, col_width, Inches(0.6))
    m2_frame = m2_box.text_frame
    m2_para = m2_frame.paragraphs[0]
    m2_para.alignment = PP_ALIGN.CENTER

    run3 = m2_para.add_run()
    run3.text = "1-3 Months"
    run3.font.size = Pt(32)
    run3.font.bold = True
    run3.font.color.rgb = white

    run4 = m2_para.add_run()
    run4.text = "\nPayback Period"
    run4.font.size = Pt(14)
    run4.font.color.rgb = white

    # Metric 3: Net Benefit
    m3_box = slide.shapes.add_textbox(roi_left + Inches(6.3), metric_top, col_width, Inches(0.6))
    m3_frame = m3_box.text_frame
    m3_para = m3_frame.paragraphs[0]
    m3_para.alignment = PP_ALIGN.CENTER

    run5 = m3_para.add_run()
    run5.text = "$1-2.1M"
    run5.font.size = Pt(32)
    run5.font.bold = True
    run5.font.color.rgb = white

    run6 = m3_para.add_run()
    run6.text = "\nNet Benefit"
    run6.font.size = Pt(14)
    run6.font.color.rgb = white

    # ===== KEY BENEFITS (4 quadrants) =====
    quad_top = Inches(3.3)
    quad_height = Inches(1.5)
    quad_width = Inches(4.3)

    # Quadrant 1: Revenue Excellence
    q1_left = Inches(0.5)
    q1_box = slide.shapes.add_textbox(q1_left, quad_top, quad_width, quad_height)
    q1_frame = q1_box.text_frame
    q1_frame.word_wrap = True

    q1_title = q1_frame.paragraphs[0]
    q1_title.text = "💰 Revenue Excellence"
    q1_title.font.size = Pt(16)
    q1_title.font.bold = True
    q1_title.font.color.rgb = dark_blue
    q1_title.space_after = Pt(6)

    bullets = [
        "Reduce days in A/R: 45-60 → 30-35 days",
        "Increase clean claim rate: 75% → 95%",
        "Improve collection rate by 15-20%",
        "Reduce billing staff time by 35%"
    ]
    for bullet in bullets:
        p = q1_frame.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(11)
        p.level = 0
        p.space_after = Pt(3)

    # Quadrant 2: Clinical Excellence
    q2_left = Inches(5.2)
    q2_box = slide.shapes.add_textbox(q2_left, quad_top, quad_width, quad_height)
    q2_frame = q2_box.text_frame
    q2_frame.word_wrap = True

    q2_title = q2_frame.paragraphs[0]
    q2_title.text = "⚕️ Clinical Excellence"
    q2_title.font.size = Pt(16)
    q2_title.font.bold = True
    q2_title.font.color.rgb = dark_blue
    q2_title.space_after = Pt(6)

    bullets2 = [
        "Increase visits per day by 20-25%",
        "Reduce documentation time by 30%",
        "Decrease medication errors by 70%",
        "Enable telehealth revenue streams"
    ]
    for bullet in bullets2:
        p = q2_frame.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(11)
        p.level = 0
        p.space_after = Pt(3)

    # Quadrant 3: Patient Engagement
    q3_left = Inches(0.5)
    q3_top = Inches(5.0)
    q3_box = slide.shapes.add_textbox(q3_left, q3_top, quad_width, quad_height)
    q3_frame = q3_box.text_frame
    q3_frame.word_wrap = True

    q3_title = q3_frame.paragraphs[0]
    q3_title.text = "👥 Patient Satisfaction"
    q3_title.font.size = Pt(16)
    q3_title.font.bold = True
    q3_title.font.color.rgb = dark_blue
    q3_title.space_after = Pt(6)

    bullets3 = [
        "Reduce no-shows: 15-30% → 5-8%",
        "Increase patient retention by 25%",
        "Improve HCAHPS scores by 15-20 pts",
        "24/7 patient self-service portal"
    ]
    for bullet in bullets3:
        p = q3_frame.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(11)
        p.level = 0
        p.space_after = Pt(3)

    # Quadrant 4: Operational Efficiency
    q4_left = Inches(5.2)
    q4_top = Inches(5.0)
    q4_box = slide.shapes.add_textbox(q4_left, q4_top, quad_width, quad_height)
    q4_frame = q4_box.text_frame
    q4_frame.word_wrap = True

    q4_title = q4_frame.paragraphs[0]
    q4_title.text = "⚙️ Operational Efficiency"
    q4_title.font.size = Pt(16)
    q4_title.font.bold = True
    q4_title.font.color.rgb = dark_blue
    q4_title.space_after = Pt(6)

    bullets4 = [
        "Consolidate 5-8 vendors into one",
        "Decrease admin staff time by 35%",
        "Improve schedule utilization to 90%",
        "HIPAA compliant • FHIR R4 ready"
    ]
    for bullet in bullets4:
        p = q4_frame.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(11)
        p.level = 0
        p.space_after = Pt(3)

    # ===== FOOTER / CALL TO ACTION =====
    footer_top = Inches(6.7)
    footer_box = slide.shapes.add_textbox(Inches(0.5), footer_top, Inches(9), Inches(0.6))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Next Step: 90-Day Pilot Program in 1-2 Departments to Validate Business Case"
    footer_para = footer_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.CENTER
    footer_para.font.size = Pt(16)
    footer_para.font.bold = True
    footer_para.font.color.rgb = green

    # Save presentation
    prs.save('MEDFLOW_EXECUTIVE_PRESENTATION.pptx')
    print("✅ Successfully created MEDFLOW_EXECUTIVE_PRESENTATION.pptx")

if __name__ == '__main__':
    create_executive_slide()
